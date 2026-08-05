# -*- coding: utf-8 -*-
"""BUSI EDA - Step 3: assemble EDA/BUSI_EDA.pptx from figures + summary."""
import os, json
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG  = os.path.join(ROOT, "figures")
ART  = os.path.join(ROOT, "artifacts")
S = json.load(open(os.path.join(ART, "summary.json"), encoding="utf-8"))

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT= RGBColor(0xC4, 0x4E, 0x52)
GREY  = RGBColor(0x55, 0x55, 0x55)
DARK  = RGBColor(0x22, 0x22, 0x22)

def txt(slide, l, t, w, h, lines, size=16, color=DARK, bold=False, align=PP_ALIGN.LEFT, bullet=False, sp=4):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        run = p.add_run(); run.text = (("•  " + ln) if bullet else ln)
        f = run.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
        f.name = "Calibri"
    return tb

def bar(slide, color=NAVY, h=0.14):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def add_fig(slide, fname, top=1.15, bottom=6.15, maxw=8.6, left_center=None):
    """place image fit into box (maxw wide, (bottom-top) tall). Returns (l,t,w,h) inches."""
    path = os.path.join(FIG, fname)
    iw, ih = Image.open(path).size
    boxw, boxh = maxw, (bottom - top)
    scale = min(boxw / (iw/96.0), boxh / (ih/96.0))  # 96 dpi baseline
    w = (iw/96.0) * scale; h = (ih/96.0) * scale
    l = left_center if left_center is not None else (13.333 - w)/2
    slide.shapes.add_picture(path, Inches(l), Inches(top), Inches(w), Inches(h))
    return l, top, w, h

def content_slide(title, fig, takeaways, tnum):
    s = prs.slides.add_slide(BLANK); bar(s)
    txt(s, 0.4, 0.28, 12.5, 0.7, title, size=26, bold=True, color=NAVY)
    # figure on left, takeaways on right
    add_fig(s, fig, top=1.2, bottom=6.9, maxw=8.5, left_center=0.35)
    txt(s, 9.15, 1.25, 3.95, 0.5, "Key takeaways", size=17, bold=True, color=ACCENT)
    txt(s, 9.15, 1.85, 3.95, 5.2, takeaways, size=14.5, bullet=True, color=DARK, sp=10)
    txt(s, 0.4, 7.06, 3, 0.4, f"EDA {tnum}", size=11, color=GREY)
    return s

# ---------------- Title slide ----------------
s = prs.slides.add_slide(BLANK)
from pptx.enum.shapes import MSO_SHAPE
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
txt(s, 0.7, 2.65, 12, 1.2, "Exploratory Data Analysis", size=44, bold=True, color=RGBColor(255,255,255))
txt(s, 0.72, 3.9, 12, 0.8, "Breast Ultrasound Images Dataset (BUSI) — 780 images, 3 classes, with lesion masks",
    size=20, color=RGBColor(0xDD,0xE6,0xF0))
txt(s, 0.7, 5.4, 12, 0.5, "COMP9444 Group Project 047 — Breast Cancer Classification & Segmentation", size=16, color=GREY)
txt(s, 0.7, 5.85, 12, 0.5, "Source: Al-Dhabyani et al. (2020), Data in Brief 28:104863", size=13, color=GREY)

# ---------------- Dataset at a glance ----------------
s = prs.slides.add_slide(BLANK); bar(s)
txt(s, 0.4, 0.28, 12.5, 0.7, "Dataset at a Glance", size=26, bold=True, color=NAVY)
facts = [
    f"Total images: {S['total_images']}  (benign {S['per_class']['benign']} / malignant {S['per_class']['malignant']} / normal {S['per_class']['normal']})",
    f"Class balance: {S['per_class_pct']['benign']}% / {S['per_class_pct']['malignant']}% / {S['per_class_pct']['normal']}%  → imbalanced",
    f"Masks: {S['images_with_mask']}/{S['total_images']} images have ≥1 lesion mask (100%); {S['multi_mask_images']} images have multiple masks (max {S['max_masks_per_image']})",
    f"Images with a real lesion: {S['images_with_lesion']} (all normal images have empty masks)",
    f"Image sizes: NOT uniform — {S['unique_sizes']} unique sizes; width {S['width_range'][0]}–{S['width_range'][1]} px, height {S['height_range'][0]}–{S['height_range'][1]} px (median {int(S['width_median'])}×{int(S['height_median'])})",
    f"Lesion area (median): malignant {S['area_ratio_median_by_class']['malignant']*100:.1f}% vs benign {S['area_ratio_median_by_class']['benign']*100:.1f}% of image",
    f"Data quality: 1 exact duplicate (labeled 2 classes), {S['near_dup_pairs']} near-duplicate pairs, {S['near_dup_images_involved']} images ({S['near_dup_images_involved']/S['total_images']*100:.0f}%) in {S['near_dup_clusters']} clusters, {S['near_dup_pairs_cross_class']} cross-class pairs",
]
txt(s, 0.7, 1.3, 12.2, 5.6, facts, size=18, bullet=True, sp=14)

# ---------------- Analysis slides ----------------
content_slide("1. Class Distribution", "fig01_class_distribution.png", [
    "780 images, heavily imbalanced: benign 56% > malignant 27% > normal 17%.",
    "Benign is ~3.3× the normal class → majority-class bias risk.",
    "Mitigate with stratified splits, class weights, and balanced augmentation.",
    "Normal images carry no lesion → segmentation effectively uses 647 lesion images.",
], "analysis 1")

content_slide("2. Image Size & Aspect Ratio", "fig02_image_size.png", [
    f"Sizes are NOT the advertised uniform 500×500: {S['unique_sizes']} distinct sizes.",
    f"Width 190–1048 px, height 310–719 px; most images are landscape (W/H>1).",
    "A fixed input size is required — resize (e.g. 256×256).",
    "Prefer resize-with-padding to avoid distorting lesion shape/aspect.",
], "analysis 2")

content_slide("3. Sample Visualization", "fig03_samples.png", [
    "Benign: oval, well-circumscribed, hypoechoic lesions.",
    "Malignant: irregular, ill-defined, often with posterior shadowing.",
    "Burned-in text & calipers ('RT LOQ', 'RIGHT BREAST') appear on images.",
    "→ risk of shortcut learning; consider cropping/inpainting annotations.",
], "analysis 3")

content_slide("4. Mask Availability & Image–Mask Matching", "fig04_mask_matching.png", [
    "Every image has at least one mask — no orphan images or masks.",
    f"{S['multi_mask_images']} images have multiple lesion masks (16 benign, 1 malignant), up to {S['max_masks_per_image']}.",
    "For binary segmentation, merge multiple masks via pixel-wise union.",
    "Normal masks are all-zero (background only).",
], "analysis 4")

content_slide("5. Lesion Area Ratio", "fig05_area_ratio.png", [
    f"Malignant lesions are much larger: median {S['area_ratio_median_by_class']['malignant']*100:.1f}% of image vs benign {S['area_ratio_median_by_class']['benign']*100:.1f}%.",
    "Many benign lesions are tiny → strong foreground/background imbalance at pixel level.",
    "Dice/IoU are sensitive to small masks — consider Dice+BCE or Tversky loss.",
    "Lesion size is itself a discriminative cue for classification.",
], "analysis 5")

content_slide("6. Lesion Location Heatmap", "fig06_location_heatmap.png", [
    "Lesions concentrate slightly above the image centre for both classes.",
    "Sonographers centre the probe on the lesion → spatial acquisition bias.",
    "Models can exploit position; random crops/translations improve robustness.",
    "Normal class has no lesion signal (empty heatmap).",
], "analysis 6")

content_slide("7. Pixel Intensity / Contrast", "fig07_intensity.png", [
    "Mean brightness and contrast overlap heavily across the three classes.",
    "Global intensity alone is a weak class separator → need texture & shape.",
    "Wide brightness range across images → apply per-image normalization.",
    "Optional CLAHE can enhance low-contrast speckled regions.",
], "analysis 7")

content_slide("8. Lesion Shape Analysis", "fig08_shape.png", [
    "Malignant lesions are less circular, with lower solidity/extent.",
    "Higher eccentricity → more elongated / irregular boundaries.",
    "Shape descriptors are discriminative — supports boundary-aware models.",
    "Motivates Attention U-Net and boundary/contour-aware losses.",
], "analysis 8")

content_slide("9. Duplicate / Near-Duplicate Detection", "fig09_duplicates.png", [
    "1 EXACT duplicate (identical file) is labeled BOTH benign and malignant.",
    f"{S['near_dup_pairs']} near-duplicate pairs (pHash≤5); {S['near_dup_images_involved']} images ({S['near_dup_images_involved']/780*100:.0f}%) in {S['near_dup_clusters']} clusters.",
    f"{S['near_dup_pairs_cross_class']} near-duplicate pairs cross class boundaries (label conflicts).",
    "Known BUSI issue (Pawłowska et al., Letter to the Editor, 2023).",
], "analysis 9")

content_slide("10. Train–Test Leakage Check", "fig10_leakage.png", [
    "A naive random 70/30 split leaks ~59 near-duplicate clusters across train/test.",
    "Duplicates on both sides inflate reported accuracy/Dice (optimistic bias).",
    "Use a GROUPED split: keep each duplicate cluster (ideally each patient) on one side.",
    "De-duplicate and resolve cross-class conflicts before training.",
], "analysis 10")

# ---------------- Data-quality summary ----------------
s = prs.slides.add_slide(BLANK); bar(s, ACCENT)
txt(s, 0.4, 0.28, 12.5, 0.7, "Data-Quality Issues & Recommendations", size=26, bold=True, color=ACCENT)
txt(s, 0.6, 1.2, 6.1, 0.5, "Issues found", size=18, bold=True, color=NAVY)
txt(s, 0.6, 1.75, 6.1, 5.2, [
    "Class imbalance (56 / 27 / 17%).",
    "Non-uniform image sizes (639 unique).",
    "Burned-in annotations, calipers & text.",
    "1 exact + 186 near-duplicate pairs; 35% of images affected.",
    "10 cross-class duplicate pairs (label conflicts).",
    "Center/acquisition spatial bias.",
    "Tiny benign lesions → pixel imbalance.",
], size=15.5, bullet=True, sp=10)
txt(s, 7.0, 1.2, 6.0, 0.5, "Recommended actions", size=18, bold=True, color=NAVY)
txt(s, 7.0, 1.75, 6.0, 5.2, [
    "De-duplicate; manually resolve cross-class conflicts.",
    "Grouped / patient-aware train-val-test split.",
    "Resize-with-padding to 256×256; per-image normalization.",
    "Stratified split + class weights + balanced augmentation.",
    "Merge multi-masks (union) for binary segmentation.",
    "Crop/inpaint burned-in annotations to avoid shortcuts.",
    "Report per-class metrics; Dice+BCE for small lesions.",
], size=15.5, bullet=True, sp=10)

# ---------------- Modeling implications ----------------
s = prs.slides.add_slide(BLANK); bar(s)
txt(s, 0.4, 0.28, 12.5, 0.7, "Implications for Modeling (Project 047)", size=26, bold=True, color=NAVY)
txt(s, 0.6, 1.25, 12.3, 5.6, [
    "Classification: transfer-learned ResNet / DenseNet / EfficientNet-B0 with class weights; add Grad-CAM for interpretability.",
    "Segmentation: U-Net baseline, then DeepLabv3+ and Attention U-Net; train on benign+malignant (647 imgs); merge masks.",
    "Loss: combine Dice + BCE (or Tversky) to handle small benign lesions and pixel imbalance.",
    "End-to-end pipeline (main method): segment lesion → crop ROI → classify (Bruno et al. 2025 style), reusing the same backbones.",
    "Evaluation: grouped cross-validation with NO duplicate leakage; report per-class Accuracy/F1/AUC and Dice/IoU.",
    "Preprocessing settled by EDA: resize+pad 256×256, normalize, augment (flip/rotate/scale), optional CLAHE.",
], size=17, bullet=True, sp=14)

out = os.path.join(ROOT, "BUSI_EDA.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
